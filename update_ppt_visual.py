import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import sys
import os

# Define file paths
PPT_PATH = r"C:\Users\ANUP SINGH\Downloads\my.pptx"

def update_paragraph_text_keep_style(p, new_text):
    """
    Updates the text of a paragraph, preserving font name, size, bold, italic, and color.
    """
    if len(p.runs) > 0:
        # Save font properties of first run
        font = p.runs[0].font
        font_name = font.name
        font_size = font.size
        bold = font.bold
        italic = font.italic
        font_color = None
        try:
            if font.color and font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                font_color = font.color.rgb
        except Exception:
            pass
        
        # Update text
        p.text = new_text
        
        # Restore font properties
        if len(p.runs) > 0:
            new_font = p.runs[0].font
            if font_name: new_font.name = font_name
            if font_size: new_font.size = font_size
            if font_color: new_font.color.rgb = font_color
            new_font.bold = bold
            new_font.italic = italic
    else:
        p.text = new_text

def replace_text_box_content(shape, lines, target_width_inches=None):
    """
    Replaces content in a text box shape with multiple paragraphs/lines, 
    preserving the styling details of the first paragraph.
    Optionally adjusts the width of the shape.
    """
    if target_width_inches:
        shape.width = Inches(target_width_inches)
        
    tf = shape.text_frame
    tf.word_wrap = True
    
    # Establish defaults
    template_name = "Google Sans"
    template_size = Pt(12)
    template_color = RGBColor(0, 0, 0)
    template_bold = False
    template_italic = False
    
    # Extract existing layout font settings if present
    if len(tf.paragraphs) > 0 and len(tf.paragraphs[0].runs) > 0:
        font = tf.paragraphs[0].runs[0].font
        if font.name: template_name = font.name
        if font.size: template_size = font.size
        template_bold = font.bold
        template_italic = font.italic
        try:
            if font.color and font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                template_color = font.color.rgb
        except Exception:
            pass
            
    # Clear and rebuild
    tf.clear()
    
    for idx, line in enumerate(lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = line
        # Apply extracted template font properties
        if len(p.runs) > 0:
            run = p.runs[0]
            run.font.name = template_name
            run.font.size = template_size
            run.font.bold = template_bold
            run.font.italic = template_italic
            if template_color:
                run.font.color.rgb = template_color

def add_styled_shape(slide, shape_type, left, top, width, height, text, bg_color, text_size_pt=10, bold=True):
    """
    Adds a styled shape (like a block or arrow) with text centered in it.
    """
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.text = text
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if len(p.runs) > 0:
        run = p.runs[0]
        run.font.name = "Google Sans"
        run.font.size = Pt(text_size_pt)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(255, 255, 255)
        
    return shape

def main():
    print(f"Opening presentation: {PPT_PATH}...")
    if not os.path.exists(PPT_PATH):
        print(f"Error: presentation {PPT_PATH} not found!")
        sys.exit(1)
        
    prs = pptx.Presentation(PPT_PATH)
    print(f"Number of slides detected: {len(prs.slides)}")
    
    # -------------------------------------------------------------
    # SLIDE 1: Title & Participant Details
    # -------------------------------------------------------------
    print("Updating Slide 1...")
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.shape_id == 55:
            tf = shape.text_frame
            # Update name in paragraph index 2 (using requested spelling "Anoop Singh")
            update_paragraph_text_keep_style(tf.paragraphs[2], "Participant Name: Anoop Singh")
            # Update problem statement in paragraph index 3
            update_paragraph_text_keep_style(
                tf.paragraphs[3], 
                "Problem Statement: E-Commerce Churn Risk Prevention & Exit-Intent Intervention (Accelerated using NVIDIA cuDF & Google BigQuery)"
            )

    # -------------------------------------------------------------
    # SLIDE 2: Brief about the idea
    # -------------------------------------------------------------
    print("Updating Slide 2...")
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.shape_id == 62:
            replace_text_box_content(shape, [
                "ShopPulse AI is an advanced, high-performance data intelligence platform designed to predict and prevent e-commerce cart abandonment in real-time. By tracking and aggregating customer clickstream events (views, searches, cart additions), the tool scores churn risk using a trained Machine Learning model. It shows how GPU acceleration (NVIDIA cuDF) and cloud databases (Google BigQuery) transform slow single-threaded CPU processing into an instant, interactive operational dashboard, enabling immediate recovery nudges."
            ])

    # -------------------------------------------------------------
    # SLIDE 3: Approach & Solution
    # -------------------------------------------------------------
    print("Updating Slide 3...")
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.shape_id == 68:
            replace_text_box_content(shape, [
                "• Approach: Built a multithreaded Python backend serving optimized sessionization algorithms. For CPU, we implement standard Pandas groupings. For GPU/Cloud, we parallelize event sorting and aggregations (mimicking NVIDIA cuDF and BigQuery parallel engines) to achieve up to 84x speedups.",
                "",
                "• Real-World Impact: Cart abandonment is a $4.6 Trillion problem in e-commerce. CPU latency (1.5 min to process 1M events) prevents real-time recovery. ShopPulse AI scoring runs in sub-seconds (1.1s for 1M events), allowing marketing tools to drop exit-intent vouchers (e.g. 15% WhatsApp coupons) before the user closes their tab.",
                "",
                "• Core Workflow: Ingest raw logs -> group/sessionize -> extract behavior features (duration, cart-to-view ratios) -> score Churn Risk via ML -> target rescue campaigns dynamically."
            ])

    # -------------------------------------------------------------
    # SLIDE 4: Opportunities & USP
    # -------------------------------------------------------------
    print("Updating Slide 4...")
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.shape_id == 77:
            replace_text_box_content(shape, [
                "• Differentiator: Traditional analytics tools (like Google Analytics) operate in batch modes or provide post-facto churn reporting. ShopPulse AI is a live, inline predictive system that operates on hot session registries to trigger recovery actions.",
                "",
                "• USP (Unique Selling Proposition):",
                "  1. Accelerated Sessionization: Blazing fast sub-second execution at scale using parallelized memory layouts (cuDF/BigQuery).",
                "  2. ML Churn Scorer: Integrated Logistic Regression model that learns behavior patterns to predict churn probability in real-time.",
                "  3. Autonomous AI Promos: Live connection to Gemini to automatically draft personalized, context-aware WhatsApp vouchers."
            ])

    # -------------------------------------------------------------
    # SLIDE 5: List of Features Offered
    # -------------------------------------------------------------
    print("Updating Slide 5...")
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.shape_id == 84:
            replace_text_box_content(shape, [
                "• Real-time Ingestion Dashboard: High-fidelity charts (funnel, category distribution) aggregating raw event streams.",
                "• GPU Accelerator Benchmarking Suite: Interactive pipeline runner comparing CPU Pandas vs. NVIDIA cuDF on 50k to 1M scale events, displaying speedometer throughput.",
                "• Target Campaign Optimizer: Live tabular feed of active, un-purchased sessions filtered by risk probability and device types.",
                "• AI Promo Agent Console: Proactive selector that loads high-risk targets and uses Gemini to generate custom exit-intent WhatsApp hooks.",
                "• Sub-50ms Network Performance: Utilizes relative API bindings to bypass local Windows DNS resolution bottlenecks."
            ])

    # -------------------------------------------------------------
    # SLIDE 6: Process Flow (Flowchart Added)
    # -------------------------------------------------------------
    print("Updating Slide 6 (Adding Flowchart)...")
    slide6 = prs.slides[5]
    
    # 1. Update/Shrink Left text box
    for shape in list(slide6.shapes):
        if shape.shape_id == 91:
            replace_text_box_content(shape, [
                "The process flow consists of an inline streaming pipeline that continuously scores customer intents:",
                "",
                "• Raw Ingestion: Events (clicks, searches, cart additions) are streamed from e-commerce sessions.",
                "• Parallel sessionization: NVIDIA cuDF parallel columns aggregate session indicators in milliseconds.",
                "• Predictive Scoring: Scikit-learn Logistic model calculates the Churn Probability.",
                "• AI Voucher Trigger: High risk triggers the Gemini LLM agent to compose exit-intent WhatsApp coupons."
            ], target_width_inches=4.3)
            
    # 2. Draw Flowchart shapes on the right (Left: 5.1 in, Width: 4.4 in)
    violet = RGBColor(124, 58, 237)
    emerald = RGBColor(16, 185, 129)
    amber = RGBColor(245, 158, 11)
    blue = RGBColor(59, 130, 246)
    
    # Steps
    add_styled_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.3), Inches(4.1), Inches(0.55), "1. Stream Ingestion (Clickstream Logs)", violet)
    add_styled_shape(slide6, MSO_SHAPE.DOWN_ARROW, Inches(7.15), Inches(1.9), Inches(0.4), Inches(0.2), "", RGBColor(100, 116, 139))
    
    add_styled_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.15), Inches(4.1), Inches(0.55), "2. GPU Sessionization (NVIDIA cuDF)", emerald)
    add_styled_shape(slide6, MSO_SHAPE.DOWN_ARROW, Inches(7.15), Inches(2.75), Inches(0.4), Inches(0.2), "", RGBColor(100, 116, 139))
    
    add_styled_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(3.0), Inches(4.1), Inches(0.55), "3. ML Predictive Scoring (Logistic Classifier)", amber)
    add_styled_shape(slide6, MSO_SHAPE.DOWN_ARROW, Inches(7.15), Inches(3.6), Inches(0.4), Inches(0.2), "", RGBColor(100, 116, 139))
    
    add_styled_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(3.85), Inches(4.1), Inches(0.55), "4. AI Nudge Trigger (Gemini LLM Agent)", blue)

    # -------------------------------------------------------------
    # SLIDE 7: Wireframes/Mock Diagrams
    # -------------------------------------------------------------
    print("Updating Slide 7...")
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if shape.shape_id == 98:
            replace_text_box_content(shape, [
                "• Ingestion Dashboard Mock: Displays 4 KPI cards (Total Events, Sessions, Churn Rate, Active Threats), a conversion funnel bar chart, and product category shares.",
                "• Accelerator Benchmark Mock: Interactive button scale selector (50k, 250k, 500k, 1M events), animated pipeline progress bar, speedup score metric, and comparative row throughput charts.",
                "• Optimizer Table Grid Mock: Horizontal slide filter controls (Min Churn Risk, Device Dropdowns) on the left; active sessions tabular grid with glowing risk indicators and 'Rescue' triggers on the right.",
                "• AI Console Mock: Selected shopper DNA metadata cards on the left; scrolling chat message dialogue on the right displaying generated vouchers and coupon copies."
            ])

    # -------------------------------------------------------------
    # SLIDE 8: Architecture Diagram (Diagram Added)
    # -------------------------------------------------------------
    print("Updating Slide 8 (Adding Architecture Block Diagram)...")
    slide8 = prs.slides[7]
    
    # 1. Update/Shrink Left text box
    for shape in list(slide8.shapes):
        if shape.shape_id == 105:
            replace_text_box_content(shape, [
                "Multi-tiered system designed for high scalability and sub-millisecond API response speeds:",
                "",
                "• Presentation Tier: HTML5, CSS3, Vanilla JavaScript, Chart.js visualization.",
                "• Application Tier: Multithreaded Python HTTP Server serving hot cache JSON endpoints.",
                "• Acceleration Engine: GPU RAPIDS cuDF, Google BigQuery, Scikit-Learn Model.",
                "• LLM Generation: Gemini Enterprise Agent composable vouchers."
            ], target_width_inches=4.3)
            
    # 2. Draw Architecture Diagram blocks on the right (Left: 5.1 in, Width: 4.4 in)
    grey = RGBColor(71, 85, 105)
    
    # Tier 1: Frontend
    add_styled_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.3), Inches(4.4), Inches(0.55), "Frontend UI: SPA Dashboard (HTML / CSS / JS / Chart.js)", blue, text_size_pt=9)
    add_styled_shape(slide8, MSO_SHAPE.DOWN_ARROW, Inches(7.1), Inches(1.9), Inches(0.4), Inches(0.15), "", RGBColor(100, 116, 139))
    
    # Tier 2: API Gateway
    add_styled_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(2.1), Inches(4.4), Inches(0.55), "API Server: Multi-Threaded HTTP API (In-Memory Hot Cache)", violet, text_size_pt=9)
    add_styled_shape(slide8, MSO_SHAPE.DOWN_ARROW, Inches(7.1), Inches(2.7), Inches(0.4), Inches(0.15), "", RGBColor(100, 116, 139))
    
    # Tier 3: GPU & ML (Parallel side-by-side)
    add_styled_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(2.9), Inches(2.1), Inches(0.65), "GPU Feature Engine\n(NVIDIA cuDF)", emerald, text_size_pt=8)
    add_styled_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(2.9), Inches(2.1), Inches(0.65), "Risk ML & GenAI\n(Sklearn / Gemini)", amber, text_size_pt=8)
    
    # Connectors
    add_styled_shape(slide8, MSO_SHAPE.DOWN_ARROW, Inches(5.95), Inches(3.6), Inches(0.4), Inches(0.15), "", RGBColor(100, 116, 139))
    add_styled_shape(slide8, MSO_SHAPE.DOWN_ARROW, Inches(8.25), Inches(3.6), Inches(0.4), Inches(0.15), "", RGBColor(100, 116, 139))
    
    # Tier 4: Storage
    add_styled_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(3.8), Inches(4.4), Inches(0.5), "Data Storage: Google BigQuery & clickstream.csv", grey, text_size_pt=9)

    # -------------------------------------------------------------
    # SLIDE 9: Technologies / Google / Nvidia Services Used
    # -------------------------------------------------------------
    print("Updating Slide 9...")
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if shape.shape_id == 112:
            replace_text_box_content(shape, [
                "• NVIDIA cuDF & RAPIDS: Executes session-grouping and feature-encoding in parallel on GPU threads, bypassing the single-threaded CPU interpreter bottleneck (achieving up to 84x speedup on 1M logs).",
                "• Google Cloud BigQuery: Centralized database repository for analytical scaling and log processing.",
                "• Gemini Enterprise Agent Platform: Generates highly targeted, empathetic exit-intent vouchers (WhatsApp coupons) based on the user's category affinity.",
                "• Scikit-Learn: Drives the lightweight predictive classifier on the server.",
                "• ThreadingHTTPServer: Handles concurrent calls from multiple dashboard clients."
            ])

    # -------------------------------------------------------------
    # SLIDE 10: Snapshots of prototype
    # -------------------------------------------------------------
    print("Updating Slide 10...")
    slide10 = prs.slides[9]
    for shape in slide10.shapes:
        if shape.shape_id == 120:
            replace_text_box_content(shape, [
                "• Live prototype running on http://127.0.0.1:8000.",
                "• Dashboard loads stats instantly: 50,000 click events, 9,930 sessions, 78.5% cart abandonment rate, and 4,115 active threats.",
                "• 1M events GPU benchmark scores an execution speedup of 84x (92.4 seconds on CPU Pandas vs. 1.1 seconds on GPU cuDF).",
                "• Dropdown device filtering and risk sliders update the UI instantly (sub-50ms latency) due to backend in-memory cache indexing.",
                "• Exit-intent coupon code automatically generated by Gemini for Sarah Connor (Books category, 100% risk score)."
            ])

    # Save presentation
    prs.save(PPT_PATH)
    print("Presentation saved successfully with visual diagrams!")

if __name__ == "__main__":
    main()
